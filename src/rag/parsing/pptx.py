"""PowerPoint decks via python-pptx.

Decks are the format most likely to hold the *decision* a policy document
only describes -- the quarterly numbers, the approved discount ladder -- so
leaving them unparsed leaves a real hole in an enterprise corpus.

Two mappings do the structural work:

- **A slide is a page.** `Block.page` is the slide number, so a citation says
  "slide 4" and the graph's provenance points somewhere a human can open.
- **A slide title is a level-1 heading.** Decks are flat: every slide is a
  peer of every other, and there is no reliable sub-heading signal. Emitting
  each title at level 1 gives `build_section_tree` one section per slide,
  which is the honest structure rather than an invented hierarchy.

Body text keeps python-pptx's outline `level` as the bullet signal, because
that is the only bullet indication the format exposes reliably: whether a
paragraph renders with a glyph lives in the layout's list style, not on the
paragraph, so an indented paragraph is the best available evidence that the
author meant a list item.
"""
from __future__ import annotations

import io

from rag.parsing.base import classify_line, is_header_line
from rag.parsing.errors import MalformedDocumentError
from rag.models import Block, BlockType, ParsedDocument

# Characters authors type when the layout does not supply a bullet glyph.
_BULLET_CHARS = "•◦▪–—-*·"


def _rectangular(rows: list[list[str]]) -> list[list[str]]:
    width = max((len(r) for r in rows), default=0)
    return [r + [""] * (width - len(r)) for r in rows]


class PptxParser:
    async def parse(self, data: bytes, doc_id: str) -> ParsedDocument:
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError

        try:
            presentation = Presentation(io.BytesIO(data))
        except (PackageNotFoundError, KeyError, ValueError) as exc:
            raise MalformedDocumentError(
                f"{doc_id}: not a readable PowerPoint package ({exc})", doc_id
            ) from exc

        blocks: list[Block] = []
        header = ""
        for page, slide in enumerate(presentation.slides, start=1):
            title = slide.shapes.title
            title_text = (title.text or "").strip() if title is not None else ""
            if title_text:
                blocks.append(Block(type=BlockType.HEADING, text=title_text,
                                    page=page, level=1))
            for shape in slide.shapes:
                if title is not None and shape is title:
                    continue
                header = self._emit_shape(shape, page, blocks, header)

        return ParsedDocument(doc_id=doc_id, file_format="pptx",
                              blocks=blocks, raw_header=header)

    def _emit_shape(self, shape, page: int, blocks: list[Block], header: str) -> str:
        # Grouped shapes are containers, not content; their text lives one
        # level down and would be lost entirely without the recursion.
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                header = self._emit_shape(child, page, blocks, header)
            return header

        if getattr(shape, "has_table", False):
            rows = [[cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows]
            if rows:
                blocks.append(Block(type=BlockType.TABLE, text="", page=page,
                                    rows=_rectangular(rows)))
            return header

        if not getattr(shape, "has_text_frame", False):
            return header

        for paragraph in shape.text_frame.paragraphs:
            text = (paragraph.text or "").strip()
            if not text:
                continue
            if not header and is_header_line(text):
                header = text
                continue
            stripped = text.lstrip(_BULLET_CHARS).strip()
            if paragraph.level > 0 or stripped != text:
                blocks.append(Block(type=BlockType.LIST, text=stripped or text,
                                    page=page))
                continue
            blocks.append(classify_line(text, page))
        return header
