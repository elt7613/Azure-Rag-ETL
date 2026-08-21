"""Regenerate the parser fixtures in this directory.

The fixtures are committed rather than built at test time so that a test
failure is reproducible from the repo alone, but they are *generated* rather
than hand-authored so their provenance is auditable: run

    python tests/fixtures/build_fixtures.py

to rebuild every file byte-for-byte from the definitions below.

The scanned PDF matters most. It is a real page of rendered pixels with no
text layer at all -- produced by drawing into a Pillow image and saving that
image as a PDF -- which is exactly what a photocopied policy document looks
like to pdfplumber: zero extractable characters. That is the input the
scanned-PDF fallback exists for, and a hand-written empty PDF would not
exercise the OCR path because there would be nothing for Document
Intelligence to read.
"""
from __future__ import annotations

from pathlib import Path

HERE = Path(__file__).parent

MARKDOWN = """\
# Remote Work Policy

People Operations | Effective: 2026-01-01 | Version 3.1 | Owner: VP People

## 1 Eligibility

Employees who have completed probation may request remote work.

- Manager approval is required.
- The employee must have a suitable workspace.

## 2 Equipment Stipend

| Item | Amount |
| --- | --- |
| Desk | $400 |
| Chair | $300 |

```python
stipend = 700
```
"""

HTML = """\
<!DOCTYPE html>
<html>
<head>
  <title>Security Standard</title>
  <style>body { color: #333; }</style>
  <script>console.log("this must never reach a block");</script>
</head>
<body>
  <h1>Security Standard</h1>
  <p>Information Security | Effective: 2026-02-01 | Version 2.0</p>
  <h2>1 Password Rules</h2>
  <p>Passwords must be at least 14 characters &amp; rotated annually.</p>
  <ul>
    <li>No reuse of the previous five passwords.</li>
    <li>MFA is mandatory for administrators.</li>
  </ul>
  <h3>1.1 Lockout Thresholds</h3>
  <table>
    <tr><th>Attempts</th><th>Lockout</th></tr>
    <tr><td>5</td><td>15 minutes</td></tr>
    <tr><td>10</td><td>24 hours</td></tr>
  </table>
</body>
</html>
"""

CSV = """\
Sales Operations | Effective: 2026-01-01 | Version 4.2,,
Tier,List Price,Discount
Starter,$32,5%
Professional,$65,10%
Enterprise,$120,18%

Region,Quota
EMEA,1200000
AMER,1800000
"""

TSV = "Country\tVAT Rate\nIreland\t23%\nGermany\t19%\n"

TXT = """\
Travel Reimbursement Guide

Finance | Effective: 2026-03-01 | Version 1.4

1 Scope

This guide covers all business travel booked through the corporate agent.

1.1 Airfare

Economy class is the default for flights under six hours.
"""

JSON = """\
{
  "policy": "Expense Limits",
  "department": "finance",
  "meta": {
    "version": "2.3",
    "owner": "Controller"
  },
  "limits": [
    {"category": "Meals", "daily_cap": 75, "currency": "USD"},
    {"category": "Lodging", "daily_cap": 250, "currency": "USD"},
    {"category": "Ground", "daily_cap": 60, "currency": "USD"}
  ]
}
"""

# cp1252-only bytes (0x92 curly apostrophe, 0x93/0x94 curly quotes) that are
# invalid UTF-8. A parser that assumes UTF-8 raises UnicodeDecodeError here.
LATIN1_TEXT = (
    b"1 Notice Period\r\n\r\n"
    b"The employee\x92s notice period is \x9330 days\x94 unless stated otherwise.\r\n"
)

# A UTF-8 BOM immediately followed by an ATX heading. Decoded naively the BOM
# becomes the first character of the line, so "﻿# Title" no longer looks
# like a heading at all.
BOM_MARKDOWN = "﻿# Vendor Onboarding\n\nLegal | Effective: 2026-04-01 | Version 1.0\n"


def _write_text(name: str, text: str) -> None:
    (HERE / name).write_text(text, encoding="utf-8")


def build_pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    title_and_body = prs.slide_layouts[1]

    slide = prs.slides.add_slide(title_and_body)
    slide.shapes.title.text = "Quarterly Business Review"
    body = slide.placeholders[1].text_frame
    body.text = "Sales | Effective: 2026-01-15 | Version 1.2"
    for line in ("Revenue grew 18% year over year.",
                 "Enterprise tier drove most of the increase."):
        para = body.add_paragraph()
        para.text = line
        para.level = 1

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Discount Ladder"
    table = slide.shapes.add_table(
        3, 2, Inches(1), Inches(2), Inches(6), Inches(2)
    ).table
    rows = [("Tier", "Max Discount"), ("Starter", "5%"), ("Enterprise", "18%")]
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            table.cell(r, c).text = value

    prs.save(HERE / "sample.pptx")


def build_image() -> None:
    """A PNG that only OCR can read -- the image branch of the routing table."""
    from PIL import Image, ImageDraw, ImageFont

    image = Image.new("RGB", (900, 300), "white")
    draw = ImageDraw.Draw(image)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 44)
    except OSError:
        font = ImageFont.load_default(size=44)
    draw.text((40, 60), "SCANNED MEMO", fill="black", font=font)
    draw.text((40, 150), "Retention period: 7 years", fill="black", font=font)
    image.save(HERE / "sample.png")


def build_scanned_pdf() -> None:
    from PIL import Image, ImageDraw, ImageFont

    lines = [
        "ACME CORP - ARCHIVED NOTICE",
        "Facilities | Effective: 2019-06-01 | Version 1.0",
        "",
        "1 Building Access",
        "Badge access to the north wing is revoked after 8 PM.",
        "Contractors must be escorted at all times.",
    ]
    image = Image.new("RGB", (1275, 1650), "white")  # 150 dpi US Letter
    draw = ImageDraw.Draw(image)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 40)
    except OSError:
        font = ImageFont.load_default(size=40)
    y = 160
    for line in lines:
        draw.text((120, y), line, fill="black", font=font)
        y += 90
    image.save(HERE / "scanned.pdf", "PDF", resolution=150.0)


def main() -> None:
    _write_text("sample.md", MARKDOWN)
    _write_text("sample.html", HTML)
    _write_text("sample.csv", CSV)
    _write_text("sample.tsv", TSV)
    _write_text("sample.txt", TXT)
    _write_text("sample.json", JSON)
    (HERE / "latin1.txt").write_bytes(LATIN1_TEXT)
    (HERE / "bom.md").write_bytes(BOM_MARKDOWN.encode("utf-8"))
    (HERE / "broken.json").write_text('{"unterminated": ', encoding="utf-8")
    # A format with no parser at all, used to prove the typed-error path.
    (HERE / "mystery.xyz").write_bytes(b"\x00\x01\x02binary payload nobody claims")
    build_pptx()
    build_image()
    build_scanned_pdf()
    print(f"wrote fixtures into {HERE}")


if __name__ == "__main__":
    main()
